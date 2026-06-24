#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gerador de slides SCP_RA03 v2 — apresentacao completa (13 slides, 16:9).
Saida: docs/SCP_RA03_v2.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paleta de cores
# ---------------------------------------------------------------------------
C_BG     = RGBColor(0x0f, 0x11, 0x17)
C_MAIN   = RGBColor(0xe2, 0xe8, 0xf0)
C_SEC    = RGBColor(0x94, 0xa3, 0xb8)
C_ACCENT = RGBColor(0x38, 0xbd, 0xf8)
C_OK     = RGBColor(0x34, 0xd3, 0x99)
C_WARN   = RGBColor(0xf8, 0x71, 0x71)
C_GOLD   = RGBColor(0xfb, 0xbf, 0x24)
C_PANEL  = RGBColor(0x1e, 0x29, 0x3b)
C_PANEL2 = RGBColor(0x16, 0x21, 0x31)
C_HDRBG  = RGBColor(0x0c, 0x2a, 0x45)

SLD_W = Inches(13.33)
SLD_H = Inches(7.5)
FOOT  = "PUCPR  ·  Complexidade de Algoritmos  ·  RA03  ·  2026"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLD_W
    prs.slide_height = SLD_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    return slide


def box(slide, text, l, t, w, h,
        sz=15, bold=False, col=None, al=PP_ALIGN.LEFT,
        fn="Segoe UI", italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = al
    r   = p.add_run()
    r.text           = text
    r.font.size      = Pt(sz)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.name      = fn
    r.font.color.rgb = col if col else C_MAIN
    return txb


def multiline(slide, lines, l, t, w, h, sz=14, fn="Segoe UI", wrap=True):
    """Lines = list of (text, color, bold, align). First item sets the textbox."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    for i, item in enumerate(lines):
        text  = item[0]
        col   = item[1] if len(item) > 1 else C_MAIN
        bold  = item[2] if len(item) > 2 else False
        align = item[3] if len(item) > 3 else PP_ALIGN.LEFT
        space = item[4] if len(item) > 4 else 2
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment  = align
        p.space_before = Pt(space)
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(sz)
        r.font.bold      = bold
        r.font.name      = fn
        r.font.color.rgb = col
    return txb


def rect_shape(slide, l, t, w, h, fill, line=None, lw=0.5):
    shp  = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width     = Pt(lw)
    else:
        shp.line.fill.background()
    return shp


def title_bar(slide, title, subtitle=None):
    rect_shape(slide, Inches(0), Inches(0), SLD_W, Pt(3), C_ACCENT)
    box(slide, title,
        Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.65),
        sz=26, bold=True, col=C_ACCENT)
    if subtitle:
        box(slide, subtitle,
            Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.38),
            sz=12, col=C_SEC)


def footer_bar(slide):
    box(slide, FOOT,
        Inches(0), Inches(7.15), SLD_W, Inches(0.3),
        sz=8, col=C_SEC, al=PP_ALIGN.CENTER)


def panel(slide, l, t, w, h, col=None, line=None):
    return rect_shape(slide, l, t, w, h, col or C_PANEL, line or C_HDRBG)


def styled_table(slide, headers, rows, l, t, w,
                 col_widths=None, sz=11, row_alt=True):
    """
    Cria tabela estilizada. col_widths = fracoes (devem somar 1).
    Cada celula de rows pode ser str ou (str, RGBColor).
    """
    nrows = len(rows) + 1
    ncols = len(headers)
    total_h = Inches(0.38 * nrows)
    tbl = slide.shapes.add_table(nrows, ncols, l, t, w, total_h).table

    if col_widths:
        for i, frac in enumerate(col_widths):
            tbl.columns[i].width = int(w * frac)

    def _cell(cell, text, bg, fg, bold=False, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        p  = tf.paragraphs[0]
        p.alignment = align
        if p.runs:
            r = p.runs[0]
        else:
            r = p.add_run()
        r.text           = text
        r.font.size      = Pt(sz)
        r.font.bold      = bold
        r.font.name      = "Segoe UI"
        r.font.color.rgb = fg

    # Cabecalho
    for j, h in enumerate(headers):
        _cell(tbl.cell(0, j), h, C_HDRBG, C_ACCENT, bold=True)

    # Dados
    for i, row in enumerate(rows):
        bg = C_PANEL if (i % 2 == 0 or not row_alt) else C_PANEL2
        for j, val in enumerate(row):
            if isinstance(val, tuple):
                text, fg = val
            else:
                text, fg = str(val), C_MAIN
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            _cell(tbl.cell(i + 1, j), text, bg, fg, align=align)

    return tbl


def code_panel(slide, lines, l, t, w, h, sz=11):
    """Painel de codigo monoespaco."""
    rect_shape(slide, l, t, w, h, C_PANEL, C_HDRBG)
    txb = slide.shapes.add_textbox(
        l + Inches(0.15), t + Inches(0.1),
        w - Inches(0.3), h - Inches(0.2))
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        text = line[0] if isinstance(line, tuple) else line
        col  = line[1] if isinstance(line, tuple) else C_MAIN
        bold = line[2] if isinstance(line, tuple) and len(line) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(sz)
        r.font.name      = "Courier New"
        r.font.color.rgb = col
        r.font.bold      = bold


# ---------------------------------------------------------------------------
# Slide 1 — Capa
# ---------------------------------------------------------------------------

def slide_capa(prs):
    s = blank_slide(prs)
    # Barra lateral accent esquerda
    rect_shape(s, Inches(0), Inches(0), Inches(0.25), SLD_H, C_ACCENT)
    # Barra accent direita
    rect_shape(s, Inches(13.08), Inches(0), Inches(0.25), SLD_H, C_ACCENT)
    # Painel central
    panel(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(5.9), C_PANEL, C_HDRBG)

    box(s, "Cobertura de Combinações",
        Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.9),
        sz=38, bold=True, col=C_ACCENT, al=PP_ALIGN.CENTER)
    box(s, "Minimum Set Cover",
        Inches(0.9), Inches(2.25), Inches(11.5), Inches(0.55),
        sz=22, col=C_SEC, al=PP_ALIGN.CENTER, italic=True)

    rect_shape(s, Inches(2.5), Inches(2.85), Inches(8.33), Pt(1.5), C_ACCENT)

    box(s, "Trabalho Avaliativo  ·  RA03",
        Inches(0.9), Inches(3.0), Inches(11.5), Inches(0.45),
        sz=16, bold=True, col=C_GOLD, al=PP_ALIGN.CENTER)
    box(s, "Complexidade de Algoritmos",
        Inches(0.9), Inches(3.42), Inches(11.5), Inches(0.4),
        sz=15, col=C_MAIN, al=PP_ALIGN.CENTER)
    box(s, "Pontifícia Universidade Católica do Paraná — PUCPR",
        Inches(0.9), Inches(3.82), Inches(11.5), Inches(0.38),
        sz=14, col=C_SEC, al=PP_ALIGN.CENTER)
    box(s, "Bacharelado em Ciência da Computação",
        Inches(0.9), Inches(4.17), Inches(11.5), Inches(0.35),
        sz=13, col=C_SEC, al=PP_ALIGN.CENTER)

    rect_shape(s, Inches(2.5), Inches(4.6), Inches(8.33), Pt(1), C_PANEL2)

    multiline(s, [
        ("Aluno:", C_SEC, False, PP_ALIGN.LEFT, 0),
        ("Jafte Carneiro Fagundes da Silva", C_MAIN, True, PP_ALIGN.LEFT, 2),
        ("Professor:", C_SEC, False, PP_ALIGN.LEFT, 8),
        ("Edson Emilio Scalabrin", C_MAIN, False, PP_ALIGN.LEFT, 2),
        ("Curitiba, 2026", C_SEC, False, PP_ALIGN.LEFT, 8),
    ], Inches(3.5), Inches(4.75), Inches(6.33), Inches(1.8), sz=13)


# ---------------------------------------------------------------------------
# Slide 2 — Modelagem do Problema
# ---------------------------------------------------------------------------

def slide_modelagem(prs):
    s = blank_slide(prs)
    title_bar(s, "Modelagem do Problema", "Minimum Set Cover sobre combinações binárias")
    footer_bar(s)

    # Coluna esquerda — definicao
    multiline(s, [
        ("Universo e Conjuntos", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("U = {1, 2, ..., 25}", C_OK, False, PP_ALIGN.LEFT, 6),
        ("S₁₅ = C(25,15) = 3.268.760 subconjuntos de tamanho 15", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Sₚ = C(25,p)    p ∈ {14, 13, 12, 11}", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Objetivo Formal", C_ACCENT, True, PP_ALIGN.LEFT, 8),
        ("Para cada p, encontrar SB₁₅,ₚ ⊆ S₁₅ tal que:", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("   ∀ Y ∈ Sₚ,  ∃ X ∈ SB₁₅,ₚ :  Y ⊆ X", C_OK, True, PP_ALIGN.LEFT, 4),
        ("Instância de Minimum Set Cover — NP-difícil", C_WARN, False, PP_ALIGN.LEFT, 8),
    ], Inches(0.4), Inches(1.15), Inches(5.8), Inches(5.5), sz=13)

    # Formulacao ILP
    multiline(s, [
        ("Formulação ILP (Programação Linear Inteira)", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("LP = variáveis contínuas xᵢ ∈ [0,1]  →  relaxação convexa", C_SEC, False, PP_ALIGN.LEFT, 4),
        ("ILP = variáveis binárias xᵢ ∈ {0,1}  →  exato, mas NP-difícil", C_MAIN, False, PP_ALIGN.LEFT, 3),
    ], Inches(6.5), Inches(1.15), Inches(6.5), Inches(0.9), sz=12)

    code_panel(s, [
        ("xᵢ ∈ {0,1}   para cada Xᵢ ∈ S₁₅",  C_MAIN,   False),
        ("",                                      C_MAIN,   False),
        ("minimizar:   Σ xᵢ",                    C_ACCENT, True),
        ("",                                      C_MAIN,   False),
        ("sujeito a:   Σ_{i: Y⊆Xᵢ} xᵢ ≥ 1",  C_ACCENT, True),
        ("             para todo Y ∈ Sₚ",         C_SEC,    False),
    ], Inches(6.5), Inches(2.05), Inches(6.5), Inches(1.9), sz=12)

    # Tabela de cardinalidades
    multiline(s, [
        ("Cardinalidades", C_ACCENT, True, PP_ALIGN.LEFT, 0),
    ], Inches(6.5), Inches(4.05), Inches(6.5), Inches(0.35), sz=12)

    styled_table(s,
        ["Conjunto", "Fórmula", "Cardinalidade"],
        [
            ["S₁₅", "C(25,15)", "3.268.760"],
            ["S₁₄", "C(25,14)", "4.457.400"],
            ["S₁₃", "C(25,13)", "5.200.300"],
            ["S₁₂", "C(25,12)", "5.200.300"],
            ["S₁₁", "C(25,11)", "4.457.400"],
        ],
        Inches(6.5), Inches(4.4), Inches(6.5),
        col_widths=[0.2, 0.35, 0.45], sz=11)


# ---------------------------------------------------------------------------
# Slide 3 — Por que ILP direto é inviável
# ---------------------------------------------------------------------------

def slide_ilp_inviavel(prs):
    s = blank_slide(prs)
    title_bar(s, "Por que ILP Direto é Inviável",
              "Escala proibitiva em todas as instâncias")
    footer_bar(s)

    box(s, "Modelo ILP sobre S₁₅ completo — dimensões por valor de p:",
        Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.4),
        sz=13, col=C_SEC)

    styled_table(s,
        ["p", "Variáveis (|S₁₅|)", "Restrições (|Sₚ|)",
         "Coef. não-nulos", "Viável?"],
        [
            ["14", "3.268.760", "4.457.400",
             "49.031.400",   ("Não", C_WARN)],
            ["13", "3.268.760", "5.200.300",
             "343.219.800",  ("Não", C_WARN)],
            ["12", "3.268.760", "5.200.300",
             "1.487.285.800", ("Não", C_WARN)],
            ["11", "3.268.760", "4.457.400",
             "4.461.857.400", ("Não", C_WARN)],
        ],
        Inches(0.4), Inches(1.65), Inches(12.5),
        col_widths=[0.06, 0.22, 0.22, 0.30, 0.20], sz=12)

    multiline(s, [
        ("Coeficientes não-nulos = C(25,p) × C(25−p, 15−p)", C_SEC, False, PP_ALIGN.LEFT, 0),
        ("Exemplo p=11:  4.457.400 × 1.001 = 4,46 bilhões de coeficientes", C_MAIN, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.75), sz=12)

    # Duas caixas explicativas
    panel(s, Inches(0.4), Inches(4.35), Inches(5.9), Inches(2.4), C_PANEL, C_HDRBG)
    multiline(s, [
        ("LP — Relaxação Contínua", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("xᵢ ∈ [0,1]  →  programação linear", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("Solúvel em tempo polinomial (Simplex/IPM)", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Fornece lower bound (LB-LP)", C_OK, False, PP_ALIGN.LEFT, 4),
        ("Não garante solução inteira válida", C_WARN, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(4.55), Inches(5.4), Inches(2.0), sz=12)

    panel(s, Inches(6.6), Inches(4.35), Inches(6.35), Inches(2.4), C_PANEL, C_HDRBG)
    multiline(s, [
        ("ILP — Programação Linear Inteira", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("xᵢ ∈ {0,1}  →  variáveis binárias", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("Exato: encontra SB mínimo global", C_OK, False, PP_ALIGN.LEFT, 4),
        ("NP-difícil: inviável em escala completa", C_WARN, False, PP_ALIGN.LEFT, 4),
        ("Estratégia: ILP pós-greedy sobre candidatos reduzidos", C_SEC, False, PP_ALIGN.LEFT, 4),
    ], Inches(6.85), Inches(4.55), Inches(5.9), Inches(2.0), sz=12)


# ---------------------------------------------------------------------------
# Slide 4 — Lower Bounds
# ---------------------------------------------------------------------------

def slide_lower_bounds(prs):
    s = blank_slide(prs)
    title_bar(s, "Lower Bounds — Critério de Qualidade",
              "Limites inferiores para |SB₁₅,ₚ|")
    footer_bar(s)

    # Explicacao esquerda
    multiline(s, [
        ("LB-LP  (contagem direta)", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Cada X ∈ S₁₅ cobre exatamente C(15,p) subconjuntos de tamanho p.", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("Logo: |SB| ≥ ⌈ C(25,p) / C(15,p) ⌉", C_OK, True, PP_ALIGN.LEFT, 5),
        ("Este é o mesmo limite que a relaxação LP fornece.", C_SEC, False, PP_ALIGN.LEFT, 4),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Cota de Schönheim", C_ACCENT, True, PP_ALIGN.LEFT, 8),
        ("Limite combinatório mais apertado, derivado da teoria de", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("códigos de cobertura (covering codes).", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Sempre: LB-Sch ≥ LB-LP", C_OK, False, PP_ALIGN.LEFT, 5),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Gap vs LB-Sch  =  |SB_greedy| / LB-Sch", C_GOLD, True, PP_ALIGN.LEFT, 8),
        ("Mede o afastamento da solução greedy em relação", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("ao melhor limite inferior conhecido.", C_MAIN, False, PP_ALIGN.LEFT, 3),
    ], Inches(0.4), Inches(1.15), Inches(5.9), Inches(5.6), sz=12)

    # Tabela de limites
    box(s, "Limites inferiores por instância",
        Inches(6.5), Inches(1.15), Inches(6.5), Inches(0.4),
        sz=13, bold=True, col=C_ACCENT)

    styled_table(s,
        ["p", "C(15,p)", "LB-LP", "LB-Schönheim"],
        [
            ["14", "15",    ("297.160",  C_MAIN), ("297.172",  C_OK)],
            ["13", "105",   ("49.527",   C_MAIN), ("58.887",   C_OK)],
            ["12", "455",   ("11.430",   C_MAIN), ("13.175",   C_OK)],
            ["11", "1.365", ("3.266",    C_MAIN), ("3.370",    C_OK)],
        ],
        Inches(6.5), Inches(1.62), Inches(6.5),
        col_widths=[0.1, 0.22, 0.34, 0.34], sz=12)

    panel(s, Inches(6.5), Inches(3.5), Inches(6.5), Inches(1.6), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Interpretação prática", C_GOLD, True, PP_ALIGN.LEFT, 0),
        ("Para p=14: o greedy precisaria de pelo menos 297.160 conjuntos.", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("A cota de Schönheim eleva esse limite para 297.172.", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Qualquer solução com |SB| = 297.172 seria ótima.", C_OK, False, PP_ALIGN.LEFT, 4),
    ], Inches(6.7), Inches(3.65), Inches(6.15), Inches(1.3), sz=12)

    panel(s, Inches(6.5), Inches(5.2), Inches(6.5), Inches(1.6), C_PANEL2, C_HDRBG)
    multiline(s, [
        ("Garantia teórica do greedy", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("|SB_greedy| ≤ H(|Sₚ|) × |OPT|", C_OK, True, PP_ALIGN.LEFT, 5),
        ("H(n) ≈ ln(n) + 0,577    →    H(5.200.300) ≈ 15,4", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Garantia teórica: gap ≤ 15,4×   —   resultados: 1,79× a 3,78×", C_OK, False, PP_ALIGN.LEFT, 4),
    ], Inches(6.7), Inches(5.35), Inches(6.15), Inches(1.3), sz=12)


# ---------------------------------------------------------------------------
# Slide 5 — Estruturas de Dados
# ---------------------------------------------------------------------------

def slide_estruturas(prs):
    s = blank_slide(prs)
    title_bar(s, "Estruturas de Dados",
              "Representação eficiente de subconjuntos de U = {1..25}")
    footer_bar(s)

    # Bitmask
    panel(s, Inches(0.4), Inches(1.15), Inches(5.9), Inches(1.9), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Bitmask uint32", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Cada subconjunto de U → inteiro de 25 bits (uint32)", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("{1,3,5} → bit 0, bit 2, bit 4 → 0b10101₂ = 21", C_OK, False, PP_ALIGN.LEFT, 4),
        ("Teste Y ⊆ X em O(1):   (X & Y) == Y", C_GOLD, True, PP_ALIGN.LEFT, 5),
    ], Inches(0.65), Inches(1.3), Inches(5.4), Inches(1.6), sz=12)

    # Indices e arrays
    panel(s, Inches(0.4), Inches(3.15), Inches(5.9), Inches(2.15), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Índices e Arrays", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("s15_index : bitmask → índice em S₁₅    O(1)", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("sp_index  : bitmask → índice em Sₚ      O(1)", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("count[]   : int32[|S₁₅|]  — cobertura atual de cada X", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("nao_coberto[] : bool[|Sₚ|] — alvos pendentes", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Arrays NumPy uint32 — operações vetorizadas", C_SEC, False, PP_ALIGN.LEFT, 5),
    ], Inches(0.65), Inches(3.3), Inches(5.4), Inches(1.85), sz=12)

    # Estrategia por p
    panel(s, Inches(0.4), Inches(5.4), Inches(5.9), Inches(1.4), C_PANEL2, C_HDRBG)
    multiline(s, [
        ("Seleção de estratégia (argmax de count[])", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("p=14, p=13 → Heap lazy (poucas atualizações/iter)", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("p=12       → np.argmax  (volume médio)", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("p=11       → Atualização paralela de count[]", C_OK, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(5.55), Inches(5.4), Inches(1.1), sz=12)

    # Heap lazy + NumPy argmax
    multiline(s, [
        ("Heap Lazy (max-heap com entradas obsoletas)", C_ACCENT, True, PP_ALIGN.LEFT, 0),
    ], Inches(6.5), Inches(1.15), Inches(6.5), Inches(0.35), sz=12)
    code_panel(s, [
        ("heap = [(-count_inicial, i) for i in range(N)]", C_MAIN),
        ("heapify(heap)",                                   C_MAIN),
        ("",                                                C_MAIN),
        ("# Extrai o máximo válido:",                       C_SEC),
        ("while True:",                                     C_MAIN),
        ("    neg_cnt, i = heappop(heap)",                  C_MAIN),
        ("    if count[i] == -neg_cnt: break  # válido",   C_OK),
    ], Inches(6.5), Inches(1.55), Inches(6.5), Inches(2.0), sz=10)

    multiline(s, [
        ("np.argmax  vs  Heap Lazy — decisão automática", C_ACCENT, True, PP_ALIGN.LEFT, 0),
    ], Inches(6.5), Inches(3.65), Inches(6.5), Inches(0.35), sz=12)

    styled_table(s,
        ["p", "Updates/iter", "Estratégia"],
        [
            ["14", "165",         ("Heap lazy",          C_MAIN)],
            ["13", "6.930",       ("Heap lazy",          C_MAIN)],
            ["12", "130.130",     ("np.argmax",          C_ACCENT)],
            ["11", "1.366.365",   ("Paralelo (Pool)",    C_OK)],
        ],
        Inches(6.5), Inches(4.05), Inches(6.5),
        col_widths=[0.1, 0.4, 0.5], sz=11)

    box(s, "Limiar empírico: > 10.000 updates/iter → argmax; >> 1M → paralelo",
        Inches(6.5), Inches(5.85), Inches(6.5), Inches(0.55),
        sz=11, col=C_SEC, italic=True)


# ---------------------------------------------------------------------------
# Slide 6 — Algoritmo Greedy + Garantia
# ---------------------------------------------------------------------------

def slide_greedy(prs):
    s = blank_slide(prs)
    title_bar(s, "Algoritmo Greedy para Set Cover",
              "Aproximação gulosa com garantia teórica logarítmica")
    footer_bar(s)

    # Pseudocodigo
    multiline(s, [
        ("Pseudocódigo", C_ACCENT, True, PP_ALIGN.LEFT, 0),
    ], Inches(0.4), Inches(1.15), Inches(6.2), Inches(0.35), sz=13)

    code_panel(s, [
        ("Entrada: S₁₅, Sₚ",                                            C_SEC,    False),
        ("Pré-proc: construir s15_index, sp_index, count[], nao_coberto[]", C_SEC, False),
        ("SB ← ∅",                                                       C_MAIN,   False),
        ("",                                                              C_MAIN,   False),
        ("enquanto existir alvo não coberto:",                            C_ACCENT, True),
        ("    X* ← argmax count[X]   // máximo atual",                   C_MAIN,   False),
        ("    SB ← SB ∪ {X*};  count[X*] ← −1",                        C_MAIN,   False),
        ("",                                                              C_MAIN,   False),
        ("    para cada Y ⊆ X*, |Y| = p:",                               C_ACCENT, False),
        ("        se Y não coberto:",                                     C_MAIN,   False),
        ("            marcar Y como coberto",                             C_OK,     False),
        ("            para cada X ∋ Y:",                                  C_MAIN,   False),
        ("                count[X] ← count[X] − 1",                     C_MAIN,   False),
        ("",                                                              C_MAIN,   False),
        ("retornar SB",                                                   C_OK,     True),
    ], Inches(0.4), Inches(1.55), Inches(6.2), Inches(5.2), sz=10)

    # Garantia + complexidade
    panel(s, Inches(6.8), Inches(1.15), Inches(6.15), Inches(2.3), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Garantia de Aproximação", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("|SB_greedy| ≤ H(|Sₚ|) × |OPT|", C_OK, True, PP_ALIGN.LEFT, 6),
        ("H(n) = ln(n) + γ,   γ ≈ 0,577   (constante de Euler)", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("H(5.200.300) ≈ 15,4", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("→ greedy garante solução ≤ 15,4× o ótimo global", C_GOLD, False, PP_ALIGN.LEFT, 4),
        ("→ na prática: 1,79× a 3,78× (muito melhor)", C_OK, False, PP_ALIGN.LEFT, 4),
    ], Inches(7.05), Inches(1.3), Inches(5.7), Inches(2.0), sz=12)

    panel(s, Inches(6.8), Inches(3.6), Inches(6.15), Inches(1.35), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Complexidade de Tempo", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Θ( K × C(k,p) × C(n−p, k−p) × log N₁₅ )", C_MAIN, True, PP_ALIGN.LEFT, 5),
        ("K = |SB|,   gargalo = atualização de count[]", C_SEC, False, PP_ALIGN.LEFT, 4),
    ], Inches(7.05), Inches(3.75), Inches(5.7), Inches(1.05), sz=12)

    panel(s, Inches(6.8), Inches(5.05), Inches(6.15), Inches(0.9), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Complexidade de Espaço", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("O(N₁₅ + Nₚ)   →   O(C(25,15) + C(25,p))", C_MAIN, False, PP_ALIGN.LEFT, 5),
    ], Inches(7.05), Inches(5.2), Inches(5.7), Inches(0.65), sz=12)

    panel(s, Inches(6.8), Inches(6.05), Inches(6.15), Inches(0.8), C_PANEL2, C_HDRBG)
    multiline(s, [
        ("O greedy não garante solução ótima global", C_WARN, True, PP_ALIGN.LEFT, 0),
        ("Encontra solução aproximada de qualidade mensurada pelo gap vs LB", C_SEC, False, PP_ALIGN.LEFT, 4),
    ], Inches(7.05), Inches(6.15), Inches(5.7), Inches(0.65), sz=11)


# ---------------------------------------------------------------------------
# Slide 7 — Análise de Complexidade
# ---------------------------------------------------------------------------

def slide_complexidade(prs):
    s = blank_slide(prs)
    title_bar(s, "Análise de Complexidade",
              "Por instância — Notações Θ, O, Ω")
    footer_bar(s)

    box(s, "Notação:  N₁₅ = C(25,15) = 3.268.760 | Nₚ = C(25,p) | K = |SB_greedy| | A = C(15,p)×C(25−p,15−p)",
        Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.38),
        sz=11, col=C_SEC)

    styled_table(s,
        ["p", "C(15,p)", "C(25−p,15−p)", "A = updates/iter", "K obtido", "Estratégia", "Tempo obs."],
        [
            ["14", "15",    "11",    ("165",       C_MAIN),   ("532.555", C_WARN),  ("Heap lazy",   C_MAIN), "~18 min"],
            ["13", "105",   "66",    ("6.930",     C_MAIN),   ("128.827", C_WARN),  ("Heap lazy",   C_MAIN), "~3 h"],
            ["12", "455",   "286",   ("130.130",   C_MAIN),   ("38.100",  C_ACCENT), ("np.argmax",  C_ACCENT), "~79 min"],
            ["11", "1.365", "1.001", ("1.366.365", C_GOLD),   ("12.733",  C_OK),    ("Paralelo",   C_OK),   "~3,4 h"],
        ],
        Inches(0.4), Inches(1.65), Inches(12.5),
        col_widths=[0.06, 0.1, 0.14, 0.18, 0.14, 0.18, 0.2], sz=11)

    # Analise programa 1
    multiline(s, [
        ("Programa 1 — Geração", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Θ(C(n,p)) por conjunto — enumera todas as combinações com itertools.combinations", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Espaço: O(C(n,p)) — aloca array uint32 completo", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Sequencial: ~80s total | Paralelo (multiprocessing.Pool): ~18,7s | Speedup: ~4,3×", C_OK, False, PP_ALIGN.LEFT, 5),
    ], Inches(0.4), Inches(3.55), Inches(6.0), Inches(1.65), sz=12)

    # Analise programas 2-5
    multiline(s, [
        ("Programas 2–5 — Greedy Set Cover", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Gargalo:  atualização de count[]  (A updates por iteração)", C_WARN, True, PP_ALIGN.LEFT, 4),
        ("Pré-proc: O(N₁₅ + Nₚ)  —  construção dos índices hash", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Loop principal: Θ(K × A × log N₁₅)  com heap;  Θ(K × (A + N₁₅))  com argmax", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Espaço: O(N₁₅ + Nₚ)  ≈  O(8,5 M) uint32  ≈  34 MB por instância", C_MAIN, False, PP_ALIGN.LEFT, 4),
    ], Inches(6.4), Inches(3.55), Inches(6.6), Inches(1.65), sz=12)

    # Comparacao
    panel(s, Inches(0.4), Inches(5.35), Inches(12.5), Inches(1.45), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Comparação: Ω, Θ, O por caso de uso", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Θ (caso médio/esperado para distribuição uniforme de Sₚ)  —  "
         "Ω (melhor caso: SB cobre tudo na 1ª iteração — improvável)  —  "
         "O (pior caso: greedy precisa de |Sₚ| iterações)", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Na prática: K << |Sₚ|;  para p=11: K=12.733 vs Nₚ=4.457.400 (0,3% das iterações do pior caso)", C_OK, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(5.5), Inches(12.0), Inches(1.15), sz=12)


# ---------------------------------------------------------------------------
# Slide 8 — Gargalos e Paralelismo
# ---------------------------------------------------------------------------

def slide_paralelismo(prs):
    s = blank_slide(prs)
    title_bar(s, "Gargalos Computacionais e Paralelismo",
              "Onde o tempo é gasto e onde o paralelismo compensa")
    footer_bar(s)

    # Programa 1
    panel(s, Inches(0.4), Inches(1.15), Inches(5.85), Inches(2.45), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Programa 1 — Geração Paralela", C_OK, True, PP_ALIGN.LEFT, 0),
        ("5 conjuntos (S₁₁..S₁₅) são independentes entre si", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("multiprocessing.Pool.map → 5 workers simultâneos", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Sequencial: ~80s   |   Paralelo: ~18,7s", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Speedup observado: ~4,3×", C_OK, True, PP_ALIGN.LEFT, 4),
        ("Overhead Windows (spawn): mais lento que Linux (fork)", C_SEC, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(1.3), Inches(5.35), Inches(2.15), sz=12)

    # Programa 5
    panel(s, Inches(0.4), Inches(3.75), Inches(5.85), Inches(2.2), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Programa 5 (p=11) — Atualização Paralela", C_OK, True, PP_ALIGN.LEFT, 0),
        ("1.366.365 updates/iter → volume justifica overhead IPC", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("Estratégia: Pool criado uma vez antes do loop", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Workers decrementam count[] em batches de Y's", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("Argmax permanece sequencial (depende do count[] completo)", C_SEC, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(3.9), Inches(5.35), Inches(1.9), sz=12)

    # Por que 2,3,4 nao paralelizados
    panel(s, Inches(6.5), Inches(1.15), Inches(6.5), Inches(4.8), C_PANEL2, C_HDRBG)
    multiline(s, [
        ("Por que p=14, 13, 12 não foram paralelizados?", C_WARN, True, PP_ALIGN.LEFT, 0),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 2),
        ("p=14 →   165 updates/iter", C_MAIN, True, PP_ALIGN.LEFT, 4),
        ("  IPC overhead >> trabalho útil por iteração", C_SEC, False, PP_ALIGN.LEFT, 3),
        ("  Heap lazy já é eficiente: log(3,27M) ≈ 22 operações/extração", C_SEC, False, PP_ALIGN.LEFT, 3),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("p=13 →   6.930 updates/iter", C_MAIN, True, PP_ALIGN.LEFT, 4),
        ("  Volume ainda pequeno vs. custo de serialização IPC", C_SEC, False, PP_ALIGN.LEFT, 3),
        ("  Heap lazy ainda compensa", C_SEC, False, PP_ALIGN.LEFT, 3),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("p=12 →   130.130 updates/iter", C_MAIN, True, PP_ALIGN.LEFT, 4),
        ("  np.argmax vetorizado: O(N₁₅) ≈ 3,27M operações NumPy/iter", C_SEC, False, PP_ALIGN.LEFT, 3),
        ("  Borderline: NumPy puro + overhead IPC = empate ou ganho marginal", C_SEC, False, PP_ALIGN.LEFT, 3),
        ("", C_MAIN, False, PP_ALIGN.LEFT, 4),
        ("p=11 →   1.366.365 updates/iter  ✓ paralelo", C_OK, True, PP_ALIGN.LEFT, 4),
        ("  Volume suficiente para amortizar overhead", C_OK, False, PP_ALIGN.LEFT, 3),
    ], Inches(6.75), Inches(1.3), Inches(6.0), Inches(4.5), sz=11)

    panel(s, Inches(0.4), Inches(6.05), Inches(5.85), Inches(0.8), C_PANEL2, C_HDRBG)
    box(s, "Critério geral: paralelizar quando updates/iter >> custo de IPC (~10k–100k threshold)",
        Inches(0.65), Inches(6.15), Inches(5.35), Inches(0.65),
        sz=11, col=C_GOLD)


# ---------------------------------------------------------------------------
# Slide 9 — Resultados Principais
# ---------------------------------------------------------------------------

def slide_resultados(prs):
    s = blank_slide(prs)
    title_bar(s, "Resultados Principais  (n=25, k=15)",
              "Greedy Set Cover — soluções obtidas vs lower bounds")
    footer_bar(s)

    styled_table(s,
        ["p", "|Sₚ|", "LB-LP", "LB-Sch", "|SB_greedy|",
         "Gap vs LB-Sch", "Tempo", "% de S₁₅"],
        [
            ["14", "4.457.400", "297.160", "297.172",
             ("532.555", C_WARN),  ("1,79×", C_GOLD), "~18 min", "16,3%"],
            ["13", "5.200.300", "49.527",  "58.887",
             ("128.827", C_WARN),  ("2,19×", C_GOLD), "~3 h",    "3,9%"],
            ["12", "5.200.300", "11.430",  "13.175",
             ("38.100",  C_ACCENT), ("2,89×", C_GOLD), "~79 min", "1,2%"],
            ["11", "4.457.400", "3.266",   "3.370",
             ("12.733",  C_OK),    ("3,78×", C_GOLD), "~3,4 h",  "0,4%"],
        ],
        Inches(0.4), Inches(1.65), Inches(12.5),
        col_widths=[0.05, 0.14, 0.13, 0.13, 0.15, 0.15, 0.12, 0.13], sz=11)

    # Mensagem principal
    panel(s, Inches(0.4), Inches(3.45), Inches(12.5), Inches(1.65), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Interpretação dos Resultados", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Garantia teórica: gap ≤ H(5.200.300) ≈ 15,4×   —   gap real: 1,79× a 3,78×", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("Greedy ficou 4× a 9× melhor que o pior caso teórico garantido", C_OK, True, PP_ALIGN.LEFT, 4),
        ("p=14 usa 16,3% de S₁₅ para cobrir 100% de S₁₄   —   p=11 usa apenas 0,4%", C_MAIN, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(3.6), Inches(12.0), Inches(1.35), sz=12)

    # Auditoria
    panel(s, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.65), C_PANEL2, C_HDRBG)
    multiline(s, [
        ("Auditoria de Cobertura Exata  (audit_cobertura.py)", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("p=11: APROVADO — 100% de S₁₁ coberto (verificação exata, sem amostragem)", C_OK, False, PP_ALIGN.LEFT, 5),
        ("p=12: APROVADO — 100% de S₁₂ coberto (verificação exata)", C_OK, False, PP_ALIGN.LEFT, 3),
        ("p=13: APROVADO — 100% de S₁₃ coberto (669,9B ops em ~80 min)", C_OK, False, PP_ALIGN.LEFT, 3),
        ("p=14: auditoria exata em andamento (2.373,8B ops  ·  ETA ~3,5 h)", C_WARN, False, PP_ALIGN.LEFT, 3),
    ], Inches(0.65), Inches(5.35), Inches(12.0), Inches(1.35), sz=12)


# ---------------------------------------------------------------------------
# Slide 10 — Eixo Comparativo das Abordagens
# ---------------------------------------------------------------------------

def slide_abordagens(prs):
    s = blank_slide(prs)
    title_bar(s, "Eixo Comparativo das Abordagens",
              "Discussão das alternativas indicadas no enunciado")
    footer_bar(s)

    styled_table(s,
        ["Abordagem", "Garantia", "Escala n=25", "Status"],
        [
            ["Greedy clássico",
             "H(Nₚ)×OPT ≈ 15,4×",
             ("Implementado e executado", C_OK),
             ("Principal ✓", C_OK)],
            ["Stochastic Greedy",
             "Sem garantia formal",
             ("Implementado — 13,5–20,6× mais rápido", C_OK),
             ("Experimental ✓", C_OK)],
            ["GRASP",
             "Sem garantia formal",
             ("Implementado — ~5× mais lento que greedy", C_MAIN),
             ("Experimental ✓", C_OK)],
            ["Algoritmos Aleatórios",
             "Esperança ≤ H(Nₚ)×OPT",
             "Viável",
             ("Alternativa futura", C_SEC)],
            ["ILP (global)",
             "Ótimo exato",
             ("Inviável em n=25 completo", C_WARN),
             ("Inviável ✗", C_WARN)],
            ["ILP pós-greedy",
             "Ótimo sobre candidatos",
             "Viável (refinamento)",
             ("Solver disponível", C_SEC)],
            ["Branch and Bound / B&P",
             "Ótimo (potencial)",
             ("Depende de heurísticas", C_MAIN),
             ("Melhoria futura", C_SEC)],
            ["Relaxação Lagrangiana",
             "Lower bound+solução",
             ("Implementado — bound 59,58 (medium)", C_OK),
             ("Experimental ✓", C_OK)],
            ["Column Generation",
             "LP exact + rounding",
             ("Implementado (medium)", C_OK),
             ("Experimental ✓", C_OK)],
            ["Computação Paralela",
             "Speedup × p",
             ("Prog.1: 4,3×; Prog.5 (p=11): paralelo", C_OK),
             ("Implementado ✓", C_OK)],
            ["Metaheurísticas (SA, GA)",
             "Sem garantia formal",
             "Viável",
             ("Melhoria futura", C_SEC)],
        ],
        Inches(0.4), Inches(1.25), Inches(12.5),
        col_widths=[0.23, 0.24, 0.33, 0.20], sz=10)


# ---------------------------------------------------------------------------
# Slide 11 — Experimentos Alternativos
# ---------------------------------------------------------------------------

def slide_experimentos(prs):
    s = blank_slide(prs)
    title_bar(s, "Experimentos Alternativos — Ganhos Medidos",
              "Benchmarks em escala reduzida para comparação de estratégias")
    footer_bar(s)

    box(s, "Instância medium  (n=15, k=9, p=6)",
        Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.38),
        sz=13, bold=True, col=C_ACCENT)

    styled_table(s,
        ["Método", "Tempo", "|SB|", "Observação"],
        [
            ["Greedy baseline",
             "10,035 s", "140",
             "Referência determinística"],
            ["Stochastic Greedy",
             ("0,743 s", C_OK), ("152", C_MAIN),
             ("13,5× mais rápido  |  +8,6% em |SB|", C_OK)],
            ["GRASP",
             "48,813 s", ("143", C_OK),
             "5× mais caro; qualidade próxima do greedy"],
            ["Relaxação Lagrangiana",
             "23,948 s", "140",
             "Bound LP: 59,58; reparo igual ao greedy"],
            ["Column Generation",
             "18,551 s", "171",
             "Relaxação LP + pricing; solução arredondada"],
        ],
        Inches(0.4), Inches(1.62), Inches(12.5),
        col_widths=[0.22, 0.13, 0.08, 0.57], sz=11)

    box(s, "Instância large-demo  (n=16, k=10, p=6)",
        Inches(0.4), Inches(3.75), Inches(12.5), Inches(0.38),
        sz=13, bold=True, col=C_ACCENT)

    styled_table(s,
        ["Método", "Tempo", "|SB|", "Observação"],
        [
            ["Greedy baseline",
             "24,895 s", "105",
             "Referência para escala maior"],
            ["Stochastic Greedy",
             ("1,208 s", C_OK), ("115", C_MAIN),
             ("20,6× mais rápido  |  +9,5% em |SB|", C_OK)],
            ["GRASP",
             "143,862 s", ("107", C_OK),
             "5,8× mais lento; melhora marginal em |SB|"],
            ["Greedy + poda local",
             "28,268 s", "105",
             "Removeu 0 candidatos redundantes"],
        ],
        Inches(0.4), Inches(4.19), Inches(12.5),
        col_widths=[0.22, 0.13, 0.08, 0.57], sz=11)

    panel(s, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.0), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Conclusão dos Experimentos", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Stochastic Greedy: maior ganho prático — 13,5–20,6× mais rápido com ≤ 10% de aumento em |SB|", C_OK, True, PP_ALIGN.LEFT, 5),
        ("GRASP valida metaheurísticas, mas não compensou o custo adicional nesta escala", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Lagrangiana e Column Generation: abordagens científicas para bounds e relaxação de LP", C_MAIN, False, PP_ALIGN.LEFT, 3),
    ], Inches(0.65), Inches(6.05), Inches(12.0), Inches(0.75), sz=11)


# ---------------------------------------------------------------------------
# Slide 12 — Limitações e Melhorias
# ---------------------------------------------------------------------------

def slide_limitacoes(prs):
    s = blank_slide(prs)
    title_bar(s, "Limitações e Possíveis Melhorias",
              "Análise crítica da solução implementada")
    footer_bar(s)

    # Limitacoes
    panel(s, Inches(0.4), Inches(1.15), Inches(6.0), Inches(5.65), C_PANEL, C_HDRBG)
    box(s, "Limitações", Inches(0.65), Inches(1.28),
        Inches(5.5), Inches(0.4), sz=14, bold=True, col=C_WARN)
    multiline(s, [
        ("ILP direto inviável em escala completa (n=25)", C_WARN, False, PP_ALIGN.LEFT, 0),
        ("  → 4,46 bilhões de coeficientes para p=11", C_SEC, False, PP_ALIGN.LEFT, 2),
        ("Greedy não garante solução ótima global", C_WARN, False, PP_ALIGN.LEFT, 8),
        ("  → gap real de 1,79× a 3,78× vs LB-Sch", C_SEC, False, PP_ALIGN.LEFT, 2),
        ("Auditoria exata p=14 é cara: ~2.373 B operações", C_WARN, False, PP_ALIGN.LEFT, 8),
        ("  → ~3,5h de verificação por NumPy sequencial", C_SEC, False, PP_ALIGN.LEFT, 2),
        ("Escalabilidade limitada para n > 27", C_WARN, False, PP_ALIGN.LEFT, 8),
        ("  → C(27,15) ≈ 21M; C(27,13) ≈ 20M; memória e tempo explodem", C_SEC, False, PP_ALIGN.LEFT, 2),
        ("Tempo alto para p=13 (~3h) e p=11 (~3,4h)", C_WARN, False, PP_ALIGN.LEFT, 8),
        ("  → gargalo: atualização de count[] (A updates/iter)", C_SEC, False, PP_ALIGN.LEFT, 2),
        ("ILP pós-greedy disponível, mas não executado em escala", C_WARN, False, PP_ALIGN.LEFT, 8),
        ("  → seria refinamento, não resultado principal", C_SEC, False, PP_ALIGN.LEFT, 2),
    ], Inches(0.65), Inches(1.72), Inches(5.5), Inches(4.9), sz=12)

    # Melhorias
    panel(s, Inches(6.7), Inches(1.15), Inches(6.25), Inches(5.65), C_PANEL2, C_HDRBG)
    box(s, "Melhorias Futuras", Inches(6.95), Inches(1.28),
        Inches(5.75), Inches(0.4), sz=14, bold=True, col=C_OK)
    multiline(s, [
        ("Algoritmos Exatos sobre Espaço Reduzido", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("  ILP pós-greedy (candidatos do greedy como universo)", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  Branch and Bound / Branch-and-Price", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  Column Generation completo (B&P integrado)", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Heurísticas Avançadas", C_ACCENT, True, PP_ALIGN.LEFT, 8),
        ("  Relaxação Lagrangiana para lower bounds mais apertados", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  Simulated Annealing pós-greedy (troca de elementos)", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  GRASP paralelo com mais restarts e busca local", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Implementação de Alto Desempenho", C_ACCENT, True, PP_ALIGN.LEFT, 8),
        ("  Bitsets vetorizados (AVX2/AVX-512)", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  Cython ou C/C++ para os gargalos Python", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  Distribuição em múltiplos processos/máquinas (MPI)", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("Auditoria Acelerada", C_ACCENT, True, PP_ALIGN.LEFT, 8),
        ("  Paralelizar audit_cobertura.py com batches de SB", C_MAIN, False, PP_ALIGN.LEFT, 3),
        ("  GPU (CuPy) para verificação de cobertura massiva", C_MAIN, False, PP_ALIGN.LEFT, 3),
    ], Inches(6.95), Inches(1.72), Inches(5.75), Inches(4.9), sz=11)


# ---------------------------------------------------------------------------
# Slide 13 — Conclusão
# ---------------------------------------------------------------------------

def slide_conclusao(prs):
    s = blank_slide(prs)
    title_bar(s, "Conclusão Técnica",
              "Síntese dos resultados e contribuições do trabalho")
    footer_bar(s)

    multiline(s, [
        ("O Problema", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("Minimum Set Cover sobre U={1..25}: NP-difícil, instâncias com até 5,2M elementos em Sₚ", C_MAIN, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.75), sz=13)

    styled_table(s,
        ["Dimensão", "Contribuição do Trabalho"],
        [
            [("Modelagem", C_ACCENT),
             "Formulação formal + ILP + lower bounds LP e Schönheim"],
            [("Algoritmo", C_ACCENT),
             "Greedy com heap lazy / np.argmax / atualização paralela (p=11)"],
            [("Estruturas", C_ACCENT),
             "Bitmask uint32 + índices hash + count[] + nao_coberto[] — acesso O(1)"],
            [("Resultados", C_ACCENT),
             "4 instâncias resolvidas; gap 1,79×–3,78× vs LB-Sch (teórico: ≤ 15,4×)"],
            [("Paralelismo", C_ACCENT),
             "Prog.1: 4,3× speedup; Prog.5 (p=11): atualização paralela de count[]"],
            [("Experimentos", C_ACCENT),
             "Stochastic Greedy, GRASP, Lagrangiana, Column Generation em benchmarks"],
            [("Auditoria", C_ACCENT),
             "p=11, 12, 13 APROVADO — cobertura 100% confirmada; p=14 em andamento"],
        ],
        Inches(0.4), Inches(2.05), Inches(12.5),
        col_widths=[0.18, 0.82], sz=12)

    panel(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.0), C_PANEL, C_HDRBG)
    multiline(s, [
        ("Mensagem Final", C_ACCENT, True, PP_ALIGN.LEFT, 0),
        ("O greedy aproximativo, com estruturas de dados eficientes e paralelismo seletivo, "
         "produziu soluções de alta qualidade (gap 1,79×–3,78×) em tempo razoável para "
         "um problema NP-difícil de escala industrial.", C_MAIN, False, PP_ALIGN.LEFT, 5),
        ("Stochastic Greedy emergiu como alternativa prática de alto desempenho: 13–20× mais "
         "rápido com custo inferior a 10% de qualidade.", C_OK, False, PP_ALIGN.LEFT, 4),
    ], Inches(0.65), Inches(6.0), Inches(12.0), Inches(0.75), sz=12)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    out_dir  = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "SCP_RA03_v2.pptx")

    prs = new_prs()

    print("Gerando slides...")
    slide_capa(prs);        print("  01/13 Capa")
    slide_modelagem(prs);   print("  02/13 Modelagem")
    slide_ilp_inviavel(prs);print("  03/13 ILP Inviável")
    slide_lower_bounds(prs);print("  04/13 Lower Bounds")
    slide_estruturas(prs);  print("  05/13 Estruturas de Dados")
    slide_greedy(prs);      print("  06/13 Algoritmo Greedy")
    slide_complexidade(prs);print("  07/13 Complexidade")
    slide_paralelismo(prs); print("  08/13 Gargalos e Paralelismo")
    slide_resultados(prs);  print("  09/13 Resultados")
    slide_abordagens(prs);  print("  10/13 Eixo Comparativo")
    slide_experimentos(prs);print("  11/13 Experimentos")
    slide_limitacoes(prs);  print("  12/13 Limitações e Melhorias")
    slide_conclusao(prs);   print("  13/13 Conclusão")

    prs.save(out_path)
    print(f"\nSalvo em: {out_path}")


if __name__ == "__main__":
    main()
