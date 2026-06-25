# -*- coding: utf-8 -*-
"""
Insert chart-only slides into docs/RA03_Complexidade.pptx.

The script does not modify existing slide content. It appends new slides, then
reorders only the new slide IDs so each chart appears near the related topic.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


ROOT = Path(__file__).resolve().parents[3]
PPTX_PATH = ROOT / "docs" / "RA03_Complexidade.pptx"
CHART_DIR = ROOT / "docs" / "assets" / "charts"
SPEAKER_GUIDE_PATH = ROOT / "docs" / "assets" / "charts" / "chart_speaker_guide.md"

EMU = 914400
BG = RGBColor(2, 6, 23)
TITLE = RGBColor(226, 232, 240)
MUTED = RGBColor(148, 163, 184)
CYAN = RGBColor(56, 189, 248)


def inch(value: float) -> Emu:
    return Emu(int(value * EMU))


def add_text(slide, x, y, w, h, text, *, size=16, color=TITLE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_chart_slide(prs: Presentation, *, kicker: str, title: str, image_name: str, footer_right: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            element = shape._element
            element.getparent().remove(element)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_text(slide, 1.04, 0.56, 17.70, 0.28, kicker, size=14, color=CYAN, bold=True)
    add_text(slide, 1.04, 0.88, 17.70, 0.56, title, size=28, color=TITLE, bold=True)

    image_path = CHART_DIR / image_name
    slide.shapes.add_picture(str(image_path), inch(0.86), inch(1.42), width=inch(18.3), height=inch(8.08))

    add_text(slide, 1.04, 10.56, 7.0, 0.24, "PUCPR · Complexidade de Algoritmos · RA03 · 2026", size=10.5, color=MUTED)
    add_text(slide, 13.6, 10.56, 4.8, 0.24, footer_right, size=10.5, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def move_slide_id(prs: Presentation, source_index: int, target_index: int) -> None:
    """Move slide ID from source_index to target_index, both zero-based."""
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[source_index]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(target_index, slide_id)


def write_speaker_guide() -> None:
    SPEAKER_GUIDE_PATH.write_text(
        """# Guia de fala - slides com gráficos

Use este guia como roteiro oral para os novos slides inseridos em `docs/RA03_Complexidade.pptx`.

## Gargalo computacional por iteração

O que dizer:
- Este gráfico explica por que o custo real não depende apenas do tamanho final de `SB`.
- Para `p=11`, a solução final tem menos conjuntos, mas cada iteração exige aproximadamente `1.366.365` atualizações de `count[]`.
- Essa diferença de ordens de grandeza justifica a escolha de paralelismo especificamente no Programa 5.
- A mensagem central é: `p=11` parece menor no resultado, mas é o mais pesado por iteração.

## Ganho com paralelismo no Programa 1

O que dizer:
- A geração dos conjuntos é naturalmente paralelizável, porque `S15`, `S14`, `S13`, `S12` e `S11` podem ser gerados de forma independente.
- A versão sequencial levou cerca de `80s`; a versão paralela levou cerca de `18,7s`.
- Isso representa speedup aproximado de `4,3x`.
- Esse resultado mostra ganho significativo de desempenho sem alterar a lógica matemática do problema.

## Resultado final versus lower bound

O que dizer:
- Este gráfico compara o tamanho obtido pelo greedy com o lower bound de Schönheim.
- A escala logarítmica permite comparar `p=14`, `p=13`, `p=12` e `p=11` no mesmo slide.
- O greedy não garante ótimo global, mas ficou relativamente próximo dos limites inferiores.
- O ponto importante é que o gap real ficou entre `1,79x` e `3,78x`, muito abaixo do pior caso teórico.

## Tempo por abordagem experimental

O que dizer:
- Este gráfico usa escala logarítmica porque os tempos variam de menos de 1 segundo a mais de 140 segundos.
- O Stochastic Greedy é claramente o método mais rápido nas duas instâncias testadas.
- No `medium`, ele foi `13,5x` mais rápido que o greedy baseline.
- No `large-demo`, foi `20,6x` mais rápido.
- O GRASP validou o eixo de metaheurísticas, mas teve custo alto nestes testes.

## Trade-off entre tempo e qualidade

O que dizer:
- O eixo X mostra tempo em escala logarítmica; quanto mais à esquerda, mais rápido.
- O eixo Y mostra `|SB|`; quanto mais abaixo, menor a solução.
- O melhor compromisso medido é o Stochastic Greedy: muito mais rápido, com aumento inferior a 10% no tamanho da solução.
- Este slide é a defesa principal da frase do enunciado sobre ganhos significativos de desempenho devidamente justificados.
""",
        encoding="utf-8",
    )


def main() -> None:
    prs = Presentation(str(PPTX_PATH))
    original_count = len(prs.slides)

    specs = [
        {
            "after_original_slide": 8,
            "kicker": "GARGALO COMPUTACIONAL",
            "title": "p=11 domina o custo por iteração",
            "image_name": "chart_updates_per_iteration.png",
        },
        {
            "after_original_slide": 8,
            "kicker": "COMPUTAÇÃO PARALELA",
            "title": "Paralelismo reduziu a geração dos conjuntos em 4,3×",
            "image_name": "chart_parallel_speedup.png",
        },
        {
            "after_original_slide": 9,
            "kicker": "RESULTADO FINAL",
            "title": "O greedy ficou perto dos lower bounds combinatórios",
            "image_name": "chart_final_sb_vs_lower_bound.png",
        },
        {
            "after_original_slide": 11,
            "kicker": "VALIDAÇÃO EXPERIMENTAL",
            "title": "Stochastic Greedy concentra o maior ganho de tempo",
            "image_name": "chart_experiment_runtime_log.png",
        },
        {
            "after_original_slide": 11,
            "kicker": "TRADE-OFF",
            "title": "O melhor compromisso medido fica no Stochastic Greedy",
            "image_name": "chart_experiment_tradeoff.png",
        },
    ]

    for idx, spec in enumerate(specs, start=1):
        add_chart_slide(
            prs,
            kicker=spec["kicker"],
            title=spec["title"],
            image_name=spec["image_name"],
            footer_right=f"Gráfico {idx} de {len(specs)}",
        )

    # Move only newly appended slides. Process target groups in reverse so the
    # original slide numbers remain stable while inserting.
    targets = [spec["after_original_slide"] for spec in specs]
    for new_offset, after_slide in sorted(enumerate(targets), key=lambda item: item[1], reverse=True):
        current_index = original_count + new_offset
        target_index = after_slide
        move_slide_id(prs, current_index, target_index)

    prs.save(str(PPTX_PATH))
    write_speaker_guide()
    print(f"Updated {PPTX_PATH}")
    print(f"Original slides: {original_count}; new slides: {len(specs)}; total: {len(prs.slides)}")
    print(f"Speaker guide: {SPEAKER_GUIDE_PATH}")


if __name__ == "__main__":
    main()
